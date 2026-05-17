"""
handlers/pptx.py — PPTX ingestion handler.

Contract: 90 System/_ingestion_contract.md §1, §2, §5, §9
HR-01 deliverable (S03).

Uses python-pptx to extract slide content. Renders slides as:

    ## Slide N — <title>

    <body text in markdown (paragraphs, bullets, tables)>

    ### Speaker Notes

    <notes text>

Media assets are extracted to:
    <staging_dir>/<original_stem>.pptx.assets/

Individual media files >20 MB are skipped with an inline marker:

    > [media skipped — file exceeds 20 MB size limit: <filename>]

Type written to frontmatter: 'presentation' (contract §2).
"""

from __future__ import annotations

import io
import re
from dataclasses import dataclass, field
from pathlib import Path
from typing import Optional

try:
    from pptx import Presentation
    from pptx.util import Pt
    from pptx.enum.text import PP_ALIGN
    _HAS_PPTX = True
except ImportError:
    _HAS_PPTX = False

try:
    from PIL import Image as PILImage
    _HAS_PIL = True
except ImportError:
    _HAS_PIL = False

from .text import HandlerResult, _PROVENANCE_VALUE

# ──────────────────────────────────────────────────────────────────────────────
# Constants
# ──────────────────────────────────────────────────────────────────────────────

_MEDIA_SIZE_LIMIT_BYTES = 20 * 1024 * 1024  # 20 MB

# Shape types that carry text frames (pptx placeholder types)
_TEXT_PLACEHOLDER_TYPES = {
    1,   # TITLE
    2,   # BODY
    3,   # CENTER_TITLE
    4,   # SUBTITLE
    5,   # OBJECT
    6,   # CHART
    7,   # TABLE
    13,  # CONTENT
    15,  # PICTURE
    19,  # SLIDE_NUMBER
    20,  # FOOTER
    21,  # DATE_TIME
}


def _pptx_version() -> str:
    try:
        import importlib.metadata
        return importlib.metadata.version("python-pptx")
    except Exception:
        return "unknown"


# ──────────────────────────────────────────────────────────────────────────────
# Shape text extraction
# ──────────────────────────────────────────────────────────────────────────────

def _paragraph_to_markdown(para) -> str:
    """Convert a pptx paragraph to a markdown line.

    Handles:
    - List bullets: indented with '  ' per level, prefixed with '-'
    - Alignment: right-aligned paragraphs get trailing '  ' (rare, best-effort)
    - Bold/italic runs: **bold**, *italic*, ***bold-italic***
    - Hyperlinks: [text](url)
    """
    # Collect run text with inline formatting
    parts: list[str] = []
    for run in para.runs:
        text = run.text or ""
        if not text.strip():
            parts.append(text)
            continue
        bold = run.font.bold
        italic = run.font.italic
        if bold and italic:
            text = f"***{text}***"
        elif bold:
            text = f"**{text}**"
        elif italic:
            text = f"*{text}*"

        # Hyperlink (pptx stores link on run.hyperlink.address)
        try:
            href = run.hyperlink.address
            if href:
                text = f"[{text}]({href})"
        except Exception:
            pass

        parts.append(text)

    line = "".join(parts)

    # Bullet / indented list
    level = para.level  # 0 = top, 1 = nested, etc.
    indent = "  " * level

    # Check if this paragraph uses a bullet/list style
    pPr = para._p.find(
        "{http://schemas.openxmlformats.org/drawingml/2006/main}buNone"
    )
    is_no_bullet = pPr is not None

    # Check for explicit bullet character or auto-number
    has_bullet_char = para._p.find(
        "{http://schemas.openxmlformats.org/drawingml/2006/main}buChar"
    ) is not None
    has_auto_num = para._p.find(
        "{http://schemas.openxmlformats.org/drawingml/2006/main}buAutoNum"
    ) is not None

    if (has_bullet_char or has_auto_num) and not is_no_bullet:
        return f"{indent}- {line}"

    if level > 0 and not is_no_bullet:
        return f"{indent}- {line}"

    return line


def _shape_to_markdown(shape) -> list[str]:
    """Extract text from a single pptx shape into a list of markdown lines."""
    lines: list[str] = []

    if not shape.has_text_frame:
        return lines

    for para in shape.text_frame.paragraphs:
        text = "".join(run.text or "" for run in para.runs)
        if not text.strip():
            # Preserve paragraph breaks as blank lines (but not excessively)
            if lines and lines[-1] != "":
                lines.append("")
            continue
        md_line = _paragraph_to_markdown(para)
        lines.append(md_line)

    return lines


def _table_shape_to_markdown(shape) -> list[str]:
    """Convert a pptx table shape to a markdown pipe table."""
    if not shape.has_table:
        return []

    table = shape.table
    rows: list[list[str]] = []
    for row in table.rows:
        cells = []
        for cell in row.cells:
            cell_text = cell.text.strip().replace("|", "\\|").replace("\n", " ")
            cells.append(cell_text or " ")
        rows.append(cells)

    if not rows:
        return []

    ncols = max(len(r) for r in rows)
    padded = [r + [""] * (ncols - len(r)) for r in rows]
    widths = [
        max(len(padded[ri][ci]) for ri in range(len(padded)))
        for ci in range(ncols)
    ]
    widths = [max(w, 3) for w in widths]

    def fmt_row(cells: list[str]) -> str:
        return "| " + " | ".join(
            c.ljust(widths[i]) for i, c in enumerate(cells)
        ) + " |"

    result = []
    for i, row in enumerate(padded):
        result.append(fmt_row(row))
        if i == 0:
            result.append("| " + " | ".join("-" * widths[j] for j in range(ncols)) + " |")

    return result


# ──────────────────────────────────────────────────────────────────────────────
# Slide title extraction
# ──────────────────────────────────────────────────────────────────────────────

def _slide_title(slide) -> str:
    """Extract slide title from the title placeholder, or fallback to first shape text."""
    # Try title placeholder (placeholder_format.idx == 0 or type == TITLE/CENTER_TITLE)
    for shape in slide.shapes:
        if shape.has_text_frame:
            try:
                pf = shape.placeholder_format
            except (ValueError, AttributeError):
                pf = None
            if pf is not None and pf.idx in (0, 1):
                text = shape.text_frame.text.strip()
                if text:
                    return text
    # Fallback: first non-empty shape text
    for shape in slide.shapes:
        if shape.has_text_frame:
            text = shape.text_frame.text.strip()
            if text:
                # Truncate at first newline
                return text.split("\n")[0][:80]
    return ""


# ──────────────────────────────────────────────────────────────────────────────
# Shape ordering heuristic (top-left reading order)
# ──────────────────────────────────────────────────────────────────────────────

def _shape_sort_key(shape):
    """Sort shapes in reading order: top-to-bottom, then left-to-right.

    Uses EMU coordinates from shape.top / shape.left.
    Title placeholder (idx=0,1) always sorts first.
    """
    try:
        pf = shape.placeholder_format
    except (ValueError, AttributeError):
        pf = None
    if pf is not None and pf.idx in (0, 1):
        return (0, 0, 0)  # title always first
    top = getattr(shape, "top", 0) or 0
    left = getattr(shape, "left", 0) or 0
    # Bucket to 1/10 of slide height for row grouping (approx 685800 EMU = 0.75cm)
    row_bucket = top // 685800
    return (1, row_bucket, left)


# ──────────────────────────────────────────────────────────────────────────────
# Media / asset extraction
# ──────────────────────────────────────────────────────────────────────────────

def _extract_media(prs, assets_dir: Path, warnings: list[str]) -> tuple[int, int]:
    """Extract all media from the PPTX package to assets_dir.

    Returns (extracted_count, skipped_count).
    Files >20 MB are skipped with a warning entry.
    """
    extracted = 0
    skipped = 0

    try:
        # python-pptx exposes the underlying zip as prs.part.package.iter_parts()
        for part in prs.part.package.iter_parts():
            # Media parts have content types starting with image/ or audio/ or video/
            ct = part.content_type or ""
            if not (ct.startswith("image/") or ct.startswith("audio/") or ct.startswith("video/")):
                continue

            blob = part.blob
            size = len(blob)

            # Derive filename from partname (e.g. "/ppt/media/image1.png")
            partname = str(part.partname)
            fname = Path(partname).name

            if size > _MEDIA_SIZE_LIMIT_BYTES:
                size_mb = size / (1024 * 1024)
                warnings.append(
                    f"media_skipped_too_large: {fname} ({size_mb:.1f} MB > 20 MB limit)"
                )
                skipped += 1
                continue

            assets_dir.mkdir(parents=True, exist_ok=True)
            dest = assets_dir / fname
            dest.write_bytes(blob)
            extracted += 1

    except Exception as exc:
        warnings.append(f"media_extraction_warning: {exc}")

    return extracted, skipped


# ──────────────────────────────────────────────────────────────────────────────
# Slide body builder
# ──────────────────────────────────────────────────────────────────────────────

def _slide_body_markdown(
    slide,
    slide_num: int,
    title: str,
    skipped_media: list[str],
) -> list[str]:
    """Build the markdown lines for a single slide (body only, not the ## header)."""
    lines: list[str] = []

    title_placeholder_idxs: set[int] = set()

    # Find the title placeholder idx so we skip it in body iteration
    for shape in slide.shapes:
        try:
            pf = shape.placeholder_format
        except (ValueError, AttributeError):
            pf = None
        if pf is not None and pf.idx in (0, 1):
            if shape.has_text_frame and shape.text_frame.text.strip():
                title_placeholder_idxs.add(id(shape))

    sorted_shapes = sorted(slide.shapes, key=_shape_sort_key)

    for shape in sorted_shapes:
        # Skip title placeholder (already in the ## header)
        if id(shape) in title_placeholder_idxs:
            continue

        # Table shapes
        if shape.has_table:
            table_lines = _table_shape_to_markdown(shape)
            if table_lines:
                lines.append("")
                lines.extend(table_lines)
                lines.append("")
            continue

        # Text shapes
        if shape.has_text_frame:
            shape_lines = _shape_to_markdown(shape)
            if shape_lines:
                lines.extend(shape_lines)
            continue

        # Image shapes (picture placeholder or picture shape)
        shape_type = getattr(shape, "shape_type", None)
        # shape_type 13 = PICTURE
        if shape_type == 13:
            name = shape.name or f"image_{slide_num}"
            lines.append(f"![{name}]")
            continue

    # Insert skipped-media markers at end of slide body
    for fname in skipped_media:
        lines.append(f"> [media skipped — file exceeds 20 MB size limit: {fname}]")

    return lines


# ──────────────────────────────────────────────────────────────────────────────
# Main handler
# ──────────────────────────────────────────────────────────────────────────────

def handle_pptx(
    source_path: Path,
    *,
    pipeline_computed: dict,
    pipeline_ns: dict,
    assets_dir: Optional[Path] = None,
) -> HandlerResult:
    """Ingest a PPTX file.

    Parameters
    ----------
    source_path : Path
        Path to the .pptx file.
    pipeline_computed : dict
        Top-level fields the pipeline computed (title, extracted_at, …).
    pipeline_ns : dict
        The pipeline: namespace dict.
    assets_dir : Path | None
        Where to extract slide media. Defaults to
        ``source_path.parent / (source_path.stem + ".pptx.assets")``.

    Output structure per slide
    --------------------------
    ## Slide N — <title>

    <body markdown>

    ### Speaker Notes

    <notes text>
    """
    if not _HAS_PPTX:
        return HandlerResult(
            success=False,
            quarantine=True,
            quarantine_reason="missing_dependency_python_pptx",
            warnings=["python-pptx not installed"],
        )

    warnings: list[str] = []

    if assets_dir is None:
        assets_dir = source_path.parent / (source_path.stem + ".pptx.assets")

    # ── Open presentation ──────────────────────────────────────────────────
    try:
        prs = Presentation(str(source_path))
    except Exception as exc:
        return HandlerResult(
            success=False,
            quarantine=True,
            quarantine_reason="pptx_open_error",
            warnings=[f"pptx_open_error: {exc}"],
        )

    slide_count = len(prs.slides)

    # ── Extract all media assets ───────────────────────────────────────────
    extracted_media, skipped_media_count = _extract_media(prs, assets_dir, warnings)

    # Build a set of skipped filenames for per-slide markers
    skipped_filenames: list[str] = [
        w.split(": ", 1)[1].split(" (")[0]
        for w in warnings
        if w.startswith("media_skipped_too_large:")
    ]

    # ── Build body ─────────────────────────────────────────────────────────
    body_sections: list[str] = []
    notes_count = 0
    table_count = 0

    for slide_num, slide in enumerate(prs.slides, start=1):
        title = _slide_title(slide)

        # Count tables on this slide
        for shape in slide.shapes:
            if shape.has_table:
                table_count += 1

        # Slide header
        header = f"## Slide {slide_num}"
        if title:
            header += f" — {title}"
        body_sections.append(header)
        body_sections.append("")

        # Slide body
        slide_lines = _slide_body_markdown(slide, slide_num, title, skipped_filenames)
        body_sections.extend(slide_lines)

        # Speaker notes
        try:
            notes_frame = slide.notes_slide.notes_text_frame if slide.has_notes_slide else None
            notes_text = notes_frame.text.strip() if notes_frame else ""
        except Exception:
            notes_text = ""

        if notes_text:
            notes_count += 1
            body_sections.append("")
            body_sections.append("### Speaker Notes")
            body_sections.append("")
            body_sections.append(notes_text)

        body_sections.append("")
        body_sections.append("---")
        body_sections.append("")

    # Remove trailing separator
    while body_sections and body_sections[-1] in ("", "---"):
        body_sections.pop()

    body = "\n".join(body_sections)

    # ── Frontmatter ────────────────────────────────────────────────────────
    top: dict = {k: v for k, v in pipeline_computed.items()}
    top["file_type"] = "pptx"
    top.setdefault("type", "presentation")
    top["slide_count"] = slide_count
    top["page_count"] = slide_count  # contract §1: page_count required

    # ── Pipeline namespace ─────────────────────────────────────────────────
    from .text import compute_text_sha256   # CH-02
    pipeline_ns["provenance"] = _PROVENANCE_VALUE
    pipeline_ns["extractor"] = f"python-pptx@{_pptx_version()}"
    pipeline_ns["slide_count"] = slide_count
    pipeline_ns["notes_count"] = notes_count
    pipeline_ns["table_count"] = table_count
    pipeline_ns["media_extracted"] = extracted_media
    pipeline_ns["media_skipped_too_large"] = skipped_media_count
    pipeline_ns["assets_dir"] = str(assets_dir) if (extracted_media > 0 or skipped_media_count > 0) else None
    pipeline_ns["text_sha256"] = compute_text_sha256(body)   # CH-02
    pipeline_ns["warnings"] = warnings

    return HandlerResult(
        success=True,
        body=body,
        frontmatter_top=top,
        pipeline_namespace=pipeline_ns,
        warnings=warnings,
    )
